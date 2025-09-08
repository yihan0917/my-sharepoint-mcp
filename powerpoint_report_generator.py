#!/usr/bin/env python3
"""
PowerPoint Report Generator for Recruiting Analysis
Creates professional PowerPoint presentations with charts and visualizations
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import io
import asyncio
import logging
from datetime import datetime
import numpy as np
from io import BytesIO
import base64
import requests
import asyncio

# Add the project root to Python path
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

from auth.sharepoint_auth import get_auth_context
from utils.graph_client import GraphClient
from config.settings import SHAREPOINT_CONFIG

def create_recruiting_presentation():
    """Create a comprehensive recruiting analysis PowerPoint presentation"""
    
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    create_title_slide(prs)
    
    # Slide 2: Executive Summary
    create_executive_summary_slide(prs)
    
    # Slide 3: Key Performance Metrics
    create_kpi_slide(prs)
    
    # Slide 4: Recruiter Performance Analysis
    create_recruiter_performance_slide(prs)
    
    # Slide 5: Time to Hire Analysis
    create_time_to_hire_slide(prs)
    
    # Slide 6: 2023 vs 2024 Comparison
    create_comparison_slide(prs)
    
    # Slide 7: Recommendations
    create_recommendations_slide(prs)
    
    return prs

def add_blue_banner_header(slide, title_text):
    """Add a blue banner header to the top of a slide with white bold text"""
    # Create a rectangle shape for the blue banner
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0),
        top=Inches(0),
        width=Inches(13.33),
        height=Inches(0.8)
    )
    
    # Set banner fill to blue (matching HR template style)
    fill = banner.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(54, 96, 146)  # Professional blue color
    
    # Remove banner outline
    banner.line.fill.background()
    
    # Add white bold text to the banner
    text_frame = banner.text_frame
    text_frame.clear()
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_right = Inches(0.3)
    text_frame.margin_top = Inches(0.15)
    text_frame.margin_bottom = Inches(0.15)
    
    p = text_frame.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    
    # Format text: white, bold, appropriate size
    run = p.runs[0]
    run.font.name = 'Calibri'
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)  # White text
    
    return banner

def create_title_slide(prs):
    """Create title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "2023 Recruiting Performance Analysis"
    subtitle.text = "Comprehensive Review & Strategic Insights\nGenerated from SharePoint Data Analysis\nSeptember 2025"
    
    # Style the title
    title_paragraph = title.text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(44)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(31, 73, 125)
    title_paragraph.alignment = PP_ALIGN.CENTER  # Center the title
    
    # Style the subtitle
    subtitle_paragraph = subtitle.text_frame.paragraphs[0]
    subtitle_paragraph.font.size = Pt(18)
    subtitle_paragraph.font.color.rgb = RGBColor(68, 84, 106)
    subtitle_paragraph.alignment = PP_ALIGN.CENTER  # Center the subtitle

def create_executive_summary_slide(prs):
    """Create executive summary slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout to avoid conflicts with banner
    slide = prs.slides.add_slide(slide_layout)
    
    # Add blue banner header
    add_blue_banner_header(slide, "Executive Summary")
    
    # Add content text box below banner
    content = slide.shapes.add_textbox(
        left=Inches(0.5),
        top=Inches(1.0),
        width=Inches(12.33),
        height=Inches(6.0)
    )
    tf = content.text_frame
    tf.text = "Key Highlights from 2023 Recruiting Performance"
    
    # Style the main text with smaller font
    main_paragraph = tf.paragraphs[0]
    main_paragraph.font.size = Pt(16)
    main_paragraph.font.bold = True
    
    # Add bullet points with smaller font sizes
    p = tf.add_paragraph()
    p.text = "106 positions successfully filled across diverse roles and locations"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "39,694 total applications processed (684 average per position)"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "74.5 days average time to fill (improved to 44-50 days in 2024)"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Top recruiters: Karrin (volume leader) and Jenna (best conversion rate)"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "70% of positions filled through internal recruiting (cost-effective approach)"
    p.level = 1
    p.font.size = Pt(14)

def create_kpi_slide(prs):
    """Create KPI slide with key metrics"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add blue banner header
    add_blue_banner_header(slide, "Key Performance Metrics - 2023")
    
    # Add content below banner
    
    # Create KPI boxes
    kpis = [
        ("Total Positions", "106", "Roles Filled"),
        ("Applications", "39,694", "Total Received"),
        ("Avg Time to Fill", "74.5", "Days"),
        ("Conversion Rate", "0.20%", "App to Offer"),
        ("Top Recruiter", "Karrin", "653 Screens"),
        ("Best Conversion", "Jenna", "13.9% Rate")
    ]
    
    # Position KPIs in 2x3 grid
    for i, (label, value, subtitle) in enumerate(kpis):
        row = i // 3
        col = i % 3
        
        x = Inches(0.5 + col * 4.2)
        y = Inches(2 + row * 2.5)
        width = Inches(3.8)
        height = Inches(2)
        
        # Create KPI box
        kpi_box = slide.shapes.add_textbox(x, y, width, height)
        kpi_frame = kpi_box.text_frame
        
        # Add label
        kpi_frame.text = label
        label_p = kpi_frame.paragraphs[0]
        label_p.font.size = Pt(14)
        label_p.font.bold = True
        label_p.alignment = PP_ALIGN.CENTER
        
        # Add value
        value_p = kpi_frame.add_paragraph()
        value_p.text = value
        value_p.font.size = Pt(36)
        value_p.font.bold = True
        value_p.font.color.rgb = RGBColor(31, 73, 125)
        value_p.alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        subtitle_p = kpi_frame.add_paragraph()
        subtitle_p.text = subtitle
        subtitle_p.font.size = Pt(12)
        subtitle_p.alignment = PP_ALIGN.CENTER

def create_recruiter_performance_slide(prs):
    """Create recruiter performance slide with chart"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add blue banner header
    add_blue_banner_header(slide, "Recruiter Performance Analysis")
    
    # Add chart data
    chart_data = CategoryChartData()
    chart_data.categories = ['Karrin', 'Jenna', 'Melissa', 'Melissa/Karrin', 'Others']
    chart_data.add_series('Recruiter Screens', (653, 79, 1, 21, 357))
    chart_data.add_series('Offers Extended', (24, 11, 12, 1, 30))
    
    # Add chart
    x, y, cx, cy = Inches(1), Inches(2), Inches(11), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

def create_time_to_hire_slide(prs):
    """Create time to hire analysis slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add blue banner header
    add_blue_banner_header(slide, "Time to Hire Analysis")
    
    # Add content with statistics
    content = slide.shapes.add_textbox(
        left=Inches(0.5),
        top=Inches(1.0),
        width=Inches(12.33),
        height=Inches(6.0)
    )
    tf = content.text_frame
    tf.text = "2023 Time to Fill Distribution"
    
    # Add statistics
    stats = [
        "Average: 74.5 days",
        "Median: 55 days",
        "Range: 3 - 400 days",
        "25th Percentile: 36 days",
        "75th Percentile: 85.5 days"
    ]
    
    for stat in stats:
        p = tf.add_paragraph()
        p.text = stat
        p.level = 1
        p.font.size = Pt(18)
    
    # Add improvement note
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "2024 Improvement: Reduced to 44-50 days average"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)

def create_comparison_slide(prs):
    """Create 2023 vs 2024 comparison slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add blue banner header
    add_blue_banner_header(slide, "2023 vs 2024 Performance Comparison")
    
    # Create comparison table
    comparisons = [
        ("Metric", "2023 (Full Year)", "2024 YTD (July)", "Improvement"),
        ("Total Hires", "106", "35", "On Track"),
        ("Avg Time to Fill", "74.5 days", "44-50 days", "33% Faster ✓"),
        ("Data Quality", "Incomplete", "Structured", "Improved ✓"),
        ("Offer Acceptance", "Not Tracked", "64-100%", "Better Tracking ✓"),
        ("Process Goals", "None", "Defined Targets", "Goal-Oriented ✓")
    ]
    
    # Add table
    rows = len(comparisons)
    cols = len(comparisons[0])
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(11), Inches(4)).table
    
    # Populate table
    for row_idx, row_data in enumerate(comparisons):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_data
            
            # Style header row
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(31, 73, 125)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
                    paragraph.font.bold = True

def create_recommendations_slide(prs):
    """Create recommendations slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add blue banner header
    add_blue_banner_header(slide, "Strategic Recommendations")
    
    # Add content
    content = slide.shapes.add_textbox(
        left=Inches(0.5),
        top=Inches(1.0),
        width=Inches(12.33),
        height=Inches(6.0)
    )
    tf = content.text_frame
    tf.text = "Immediate Actions"
    
    # Style the main text with smaller font
    main_paragraph = tf.paragraphs[0]
    main_paragraph.font.size = Pt(16)
    main_paragraph.font.bold = True
    
    immediate_actions = [
        "Improve data quality - address 45% missing application data",
        "Standardize processes - reduce time-to-hire variance",
        "Share best practices from high-performing recruiters"
    ]
    
    for action in immediate_actions:
        p = tf.add_paragraph()
        p.text = action
        p.level = 1
        p.font.size = Pt(12)
    
    # Add strategic initiatives
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Strategic Initiatives"
    p.font.bold = True
    p.font.size = Pt(16)
    
    strategic_items = [
        "Set specific time-to-hire targets by role type",
        "Investigate roles taking >90 days to fill",
        "Implement consistent conversion rate tracking"
    ]
    
    for item in strategic_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(12)
    
    # Add long-term planning
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Long-term Planning"
    p.font.bold = True
    p.font.size = Pt(16)
    
    longterm_items = [
        "Plan for 100-120 annual hires based on 2023 data",
        "Optimize recruiter assignments based on performance",
        "Understand factors causing 400-day outliers"
    ]
    
    for item in longterm_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(12)

async def upload_to_sharepoint(prs, filename="2023_Recruiting_Analysis_Presentation.pptx"):
    """Upload PowerPoint presentation to SharePoint AI Generated Reports folder"""
    try:
        print("🔐 Authenticating with SharePoint...")
        context = await get_auth_context()
        
        # Create Graph client
        graph_client = GraphClient(context)
        
        # Get site info
        site_parts = SHAREPOINT_CONFIG["site_url"].replace("https://", "").split("/")
        domain = site_parts[0]
        site_name = site_parts[2] if len(site_parts) > 2 else "root"
        site_info = await graph_client.get_site_info(domain, site_name)
        site_id = site_info["id"]
        
        # Get document libraries
        libraries_response = await graph_client.list_document_libraries(domain, site_name)
        drives = libraries_response.get("value", [])
        
        # Find Documents library
        documents_drive = None
        for drive in drives:
            if drive["name"] == "Documents":
                documents_drive = drive
                break
        
        if not documents_drive:
            raise Exception("Documents library not found")
        
        drive_id = documents_drive["id"]
        
        print("📁 Looking for AI Generated Reports folder...")
        
        # Check if AI Generated Reports folder exists, create if not
        try:
            items_response = await graph_client.list_document_contents(site_id, drive_id, "root")
            items = items_response.get("value", [])
            
            ai_reports_folder = None
            for item in items:
                if item.get("folder") and item["name"] == "AI Generated Reports":
                    ai_reports_folder = item
                    break
            
            if not ai_reports_folder:
                print("📂 Creating AI Generated Reports folder...")
                # Create folder using Graph API
                create_folder_url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/root/children"
                folder_data = {
                    "name": "AI Generated Reports",
                    "folder": {},
                    "@microsoft.graph.conflictBehavior": "rename"
                }
                
                response = requests.post(create_folder_url, headers=context.headers, json=folder_data)
                if response.status_code == 201:
                    ai_reports_folder = response.json()
                    print("✅ AI Generated Reports folder created")
                else:
                    raise Exception(f"Failed to create folder: {response.text}")
            
            folder_id = ai_reports_folder["id"]
            
        except Exception as e:
            print(f"❌ Error accessing folder: {str(e)}")
            raise
        
        print("💾 Saving presentation to memory...")
        # Save presentation to BytesIO
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        pptx_content = pptx_buffer.getvalue()
        
        print("📤 Uploading to SharePoint...")
        # Upload file to AI Generated Reports folder
        upload_url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/items/{folder_id}:/{filename}:/content"
        
        upload_headers = context.headers.copy()
        upload_headers.pop("Content-Type", None)  # Remove content-type for binary upload
        upload_headers["Content-Type"] = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        
        response = requests.put(upload_url, headers=upload_headers, data=pptx_content)
        
        if response.status_code in [200, 201]:
            upload_result = response.json()
            print(f"✅ PowerPoint uploaded successfully to SharePoint!")
            print(f"📊 File: {filename}")
            print(f"📁 Location: AI Generated Reports folder")
            print(f"🔗 SharePoint URL: {upload_result.get('webUrl', 'N/A')}")
            return upload_result
        else:
            raise Exception(f"Upload failed: HTTP {response.status_code} - {response.text}")
            
    except Exception as e:
        print(f"❌ Error uploading to SharePoint: {str(e)}")
        raise

async def main():
    """Generate and upload the PowerPoint presentation"""
    print("Generating 2023 Recruiting Analysis PowerPoint Presentation...")
    
    try:
        # Create presentation
        prs = create_recruiting_presentation()
        
        print(f"📊 Generated {len(prs.slides)} slides with comprehensive analysis")
        print("📈 Includes: KPIs, Charts, Comparisons, and Recommendations")
        
        # Upload to SharePoint instead of saving locally
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"2023_Recruiting_Analysis_Presentation_{timestamp}.pptx"
        upload_result = await upload_to_sharepoint(prs, filename)
        
        return upload_result
        
    except Exception as e:
        print(f"❌ Error generating PowerPoint: {str(e)}")
        return None

if __name__ == "__main__":
    asyncio.run(main())
